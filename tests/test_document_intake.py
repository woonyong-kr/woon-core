from __future__ import annotations

import hashlib
import json
import os
import shutil
import threading
import time
from io import StringIO
from pathlib import Path
from typing import Any

import pytest

from woon_core.cli import run
from woon_core.errors import WoonError
from woon_core.knowledge import document_intake
from woon_core.knowledge.document_intake import (
    DOCUMENT_INTAKE_SCHEMA_VERSION,
    ingest_document_candidate,
)
from woon_core.knowledge.document_resolution import (
    audit_document_resolutions,
    resolve_document_candidate,
)
from woon_core.knowledge.source_catalog import plan_source_catalog

FIXTURE = Path(__file__).parent / "fixtures/docling/heading-table.html"


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def test_docling_converts_structured_fixture_and_replays_without_duplicates(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()

    first = ingest_document_candidate(FIXTURE, vault, source_locator="fixtures/table.html")

    assert first.status == "candidate"
    assert first.replayed is False
    assert first.candidate is not None
    candidate = vault / first.candidate
    markdown = (candidate / "document.md").read_text(encoding="utf-8")
    document = json.loads((candidate / "document.json").read_text(encoding="utf-8"))
    chunks = [
        json.loads(line)
        for line in (candidate / "chunks.jsonl").read_text(encoding="utf-8").splitlines()
    ]
    receipt = json.loads((candidate / "receipt.json").read_text(encoding="utf-8"))
    observation = json.loads((vault / first.observation).read_text(encoding="utf-8"))
    promotion = json.loads((candidate / "promotion.json").read_text(encoding="utf-8"))

    assert "# 문서 수집 검증" in markdown
    assert "| DOCX" in markdown
    assert any(item["label"] == "section_header" for item in document["texts"])
    assert len(document["tables"]) == 1
    assert chunks
    assert any("지원 우선순위" in chunk["contextualized_text"] for chunk in chunks)
    raw_markdown = (candidate / "document.raw.md").read_text(encoding="utf-8")
    quality = json.loads((candidate / "quality.json").read_text(encoding="utf-8"))
    assert receipt["promotion_state"] == "curation-required"
    assert receipt["canonical_writes"] is False
    assert receipt["version"] == DOCUMENT_INTAKE_SCHEMA_VERSION
    assert receipt["adapter"]["schema_version"] == DOCUMENT_INTAKE_SCHEMA_VERSION
    assert receipt["converter"]["packages"]["docling"] == "2.117.0"
    assert "onnxruntime" not in receipt["converter"]["packages"]
    assert "rapidocr" not in receipt["converter"]["packages"]
    assert "locator" not in receipt["source"]
    assert observation["source"]["locator"] == "fixtures/table.html"
    assert observation["candidate_id"] == first.candidate_id
    assert promotion["required_gates"] == [
        "semantic-value",
        "knowledge-search",
        "duplicate-resolution",
        "terminal-resolution-receipt",
    ]
    assert promotion["curation_contract"]["promotion_interface"] == (
        "woon knowledge document-resolve"
    )
    normalized_markdown = (
        "\n".join(line.rstrip() for line in markdown.replace("\r\n", "\n").split("\n")).strip()
        + "\n"
    )
    assert (
        promotion["curation_contract"]["input_sha256"]
        == hashlib.sha256(normalized_markdown.encode()).hexdigest()
    )
    assert raw_markdown
    assert quality["state"] == "ready-for-semantic-curation"
    assert quality["canonical_writes"] is False
    for name in (
        "document.json",
        "document.raw.md",
        "document.md",
        "chunks.jsonl",
        "promotion.json",
        "quality.json",
        "receipt.json",
    ):
        assert (candidate / name).stat().st_mode & 0o777 == 0o600
    assert (vault / first.observation).stat().st_mode & 0o777 == 0o600
    receipt_hash = _sha256(candidate / "receipt.json")

    replay = ingest_document_candidate(FIXTURE, vault, source_locator="fixtures/table.html")

    assert replay.candidate_id == first.candidate_id
    assert replay.replayed is True
    assert replay.observation == first.observation
    assert _sha256(candidate / "receipt.json") == receipt_hash
    candidates = vault / ".local/woon-knowledge/document-intake/candidates"
    assert [path.name for path in candidates.iterdir()] == [first.candidate_id]


def test_failed_conversion_is_isolated_from_canonical_and_compiler_receipt(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    canonical = vault / "wiki/existing.md"
    compiler_receipt = vault / "catalog/llm-wiki/receipts/existing.json"
    canonical.parent.mkdir(parents=True)
    compiler_receipt.parent.mkdir(parents=True)
    canonical.write_text("# Existing\n", encoding="utf-8")
    compiler_receipt.write_text('{"stable": true}\n', encoding="utf-8")
    invalid = tmp_path / "invalid.docx"
    invalid.write_bytes(b"this is not an Office document")
    before = (_sha256(canonical), _sha256(compiler_receipt))

    with pytest.raises(WoonError, match="error receipt"):
        ingest_document_candidate(invalid, vault)

    assert (_sha256(canonical), _sha256(compiler_receipt)) == before
    failure_receipts = list(
        (vault / ".local/woon-knowledge/document-intake/failures").glob("*.json")
    )
    assert len(failure_receipts) == 1
    failure_text = failure_receipts[0].read_text(encoding="utf-8")
    failure = json.loads(failure_text)
    assert failure["status"] == "failed"
    assert failure["version"] == DOCUMENT_INTAKE_SCHEMA_VERSION
    assert failure["canonical_writes"] is False
    assert str(tmp_path) not in failure_text
    assert str(Path.home()) not in failure_text
    assert not (vault / ".local/woon-knowledge/document-intake/candidates").exists()


def test_pdf_requires_explicit_prefetched_model_cache(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()
    source = tmp_path / "sample.pdf"
    source.write_bytes(b"%PDF-1.4\n%%EOF\n")

    with pytest.raises(WoonError, match="pre-fetched local model cache"):
        ingest_document_candidate(source, vault)
    with pytest.raises(WoonError, match="pre-fetched local model cache"):
        ingest_document_candidate(source, vault)

    assert not (vault / "wiki").exists()
    failure_receipts = list(
        (vault / ".local/woon-knowledge/document-intake/failures").glob("*.json")
    )
    assert len(failure_receipts) == 1
    failure = json.loads(failure_receipts[0].read_text(encoding="utf-8"))
    assert failure["status"] == "failed"
    assert failure["candidate_id"] is None
    assert failure["canonical_writes"] is False


def test_model_cache_symlink_root_fails_with_receipt(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()
    source = tmp_path / "sample.pdf"
    source.write_bytes(b"%PDF-1.4\n%%EOF\n")
    cache_target = tmp_path / "cache-target"
    cache_target.mkdir()
    (cache_target / "model.bin").write_bytes(b"model")
    cache_link = tmp_path / "cache-link"
    cache_link.symlink_to(cache_target, target_is_directory=True)

    with pytest.raises(WoonError, match="symlink root"):
        ingest_document_candidate(source, vault, model_cache=cache_link)

    failures = list((vault / ".local/woon-knowledge/document-intake/failures").glob("*.json"))
    assert len(failures) == 1


def test_tampered_candidate_fails_closed_without_replacement(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()
    first = ingest_document_candidate(FIXTURE, vault)
    assert first.candidate is not None
    candidate = vault / first.candidate
    tampered = b"tampered\n"
    (candidate / "document.md").write_bytes(tampered)

    with pytest.raises(WoonError, match="error receipt"):
        ingest_document_candidate(FIXTURE, vault)

    assert (candidate / "document.md").read_bytes() == tampered
    failures = list((vault / ".local/woon-knowledge/document-intake/failures").glob("*.json"))
    assert len(failures) == 1


def test_same_content_at_two_locators_uses_one_candidate_and_two_observations(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()

    first = ingest_document_candidate(FIXTURE, vault, source_locator="drop-a/table.html")
    second = ingest_document_candidate(FIXTURE, vault, source_locator="drop-b/table.html")

    assert second.candidate_id == first.candidate_id
    assert second.replayed is True
    assert second.observation != first.observation
    candidates = vault / ".local/woon-knowledge/document-intake/candidates"
    observations = vault / ".local/woon-knowledge/document-intake/observations" / first.candidate_id
    assert len(list(candidates.iterdir())) == 1
    assert len(list(observations.glob("*.json"))) == 2


def test_same_bytes_with_different_input_formats_do_not_share_candidate(
    tmp_path: Path,
) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()
    html = tmp_path / "same.html"
    markdown = tmp_path / "same.md"
    content = "# format identity\n"
    html.write_text(content, encoding="utf-8")
    markdown.write_text(content, encoding="utf-8")

    html_result = ingest_document_candidate(html, vault)
    markdown_result = ingest_document_candidate(markdown, vault)

    assert html_result.source_sha256 == markdown_result.source_sha256
    assert html_result.candidate_id != markdown_result.candidate_id


def test_source_change_during_snapshot_fails_with_receipt(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()
    source = tmp_path / "source.html"
    source.write_text(FIXTURE.read_text(encoding="utf-8"), encoding="utf-8")
    real_copy = shutil.copyfile

    def mutate_then_copy(origin: Path, destination: Path) -> str:
        Path(origin).write_text("<h1>changed</h1>\n", encoding="utf-8")
        return str(real_copy(origin, destination))

    monkeypatch.setattr("woon_core.knowledge.document_intake.shutil.copyfile", mutate_then_copy)

    with pytest.raises(WoonError, match="error receipt"):
        ingest_document_candidate(source, vault)

    failures = list((vault / ".local/woon-knowledge/document-intake/failures").glob("*.json"))
    assert len(failures) == 1
    failure = json.loads(failures[0].read_text(encoding="utf-8"))
    assert "source changed" in failure["error"]["message"]
    assert not (vault / ".local/woon-knowledge/document-intake/candidates").exists()


def test_offline_environment_serializes_threads_and_restores_process_state() -> None:
    original = {name: os.environ.get(name) for name in document_intake._OFFLINE_ENVIRONMENT}
    active = 0
    maximum = 0
    state_lock = threading.Lock()

    def worker() -> None:
        nonlocal active, maximum
        with document_intake._offline_environment():
            assert os.environ["HF_HUB_OFFLINE"] == "1"
            assert os.environ["TRANSFORMERS_OFFLINE"] == "1"
            with state_lock:
                active += 1
                maximum = max(maximum, active)
            time.sleep(0.02)
            with state_lock:
                active -= 1

    threads = [threading.Thread(target=worker) for _ in range(3)]
    for thread in threads:
        thread.start()
    for thread in threads:
        thread.join()

    assert maximum == 1
    assert {name: os.environ.get(name) for name in document_intake._OFFLINE_ENVIRONMENT} == original


def test_standalone_image_projection_preserves_ocr_omitted_by_markdown_and_chunks() -> None:
    class TextItem:
        def __init__(self, text: str) -> None:
            self.text = text

    class Document:
        texts = [
            TextItem("Reverse Proxy Flow"),
            TextItem("user's device (D)"),
            TextItem("origin server (F)"),
            TextItem("origin server (F)"),
        ]

    markdown, chunks = document_intake._add_image_ocr_projection(
        Document(),
        "## Reverse Proxy Flow",
        [],
    )

    assert markdown == (
        "## Reverse Proxy Flow\n\n## OCR text\n\n- user's device (D)\n- origin server (F)"
    )
    assert chunks == [
        {
            "index": 0,
            "contextualized_text": (
                "OCR text\nReverse Proxy Flow\nuser's device (D)\norigin server (F)"
            ),
            "chunk": {
                "text": "Reverse Proxy Flow\nuser's device (D)\norigin server (F)",
                "meta": {
                    "origin": "standalone-image-ocr-projection",
                    "schema_version": 1,
                },
            },
        }
    ]


def test_ocr_language_contract_is_explicit() -> None:
    assert document_intake._ocr_languages("off") == []
    assert document_intake._ocr_languages("rapidocr") == ["chinese"]
    assert document_intake._ocr_languages("ocrmac") == ["ko-KR", "en-US"]


def test_image_cleaning_removes_chrome_duplicates_and_obvious_ocr_noise() -> None:
    class Bbox:
        def __init__(self, left: float, right: float) -> None:
            self.l = left
            self.r = right
            self.b = 0.0
            self.t = 100.0

    class Prov:
        def __init__(self, left: float, right: float) -> None:
            self.bbox = Bbox(left, right)

    class Item:
        def __init__(self, text: str, label: str, left: float, right: float) -> None:
            self.text = text
            self.label = label
            self.prov = [Prov(left, right)]

    class Document:
        texts = [
            Item("AICE 소개", "page_header", 600, 700),
            Item("시험 응시", "section_header", 200, 450),
            Item("시험 응시", "text", 220, 420),
            Item("사전 환경점검 및 시험 응시 유의사항 afr", "section_header", 200, 650),
            Item("~ 시험 시작시간 30분 전부터 입실 가능합니다.", "list_item", 210, 700),
            Item("FAQ", "text", 1500, 1600),
        ]

    cleaned, quality = document_intake._clean_document_markdown(
        Document(),
        "raw",
        suffix=".png",
    )

    assert cleaned == (
        "# 시험 응시\n\n"
        "## 사전 환경점검 및 시험 응시 유의사항\n\n"
        "- 시험 시작시간 30분 전부터 입실 가능합니다.\n"
    )
    assert "AICE 소개" not in cleaned
    assert "FAQ" not in cleaned
    assert "afr" not in cleaned
    assert quality["state"] == "ready-for-semantic-curation"
    assert quality["transformations"] == {
        "exact-duplicate-removed": 1,
        "leading-glyph-removed": 1,
        "off-content-region-removed": 1,
        "page-chrome-removed": 1,
        "suspicious-trailing-token-removed": 1,
    }


def test_missing_docling_dependency_creates_failure_receipt(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()

    def unavailable(ocr: str) -> dict[str, str]:
        assert ocr == "off"
        raise WoonError("Docling is not installed")

    monkeypatch.setattr(document_intake, "_docling_versions", unavailable)

    with pytest.raises(WoonError, match="error receipt"):
        ingest_document_candidate(FIXTURE, vault)

    failures = list((vault / ".local/woon-knowledge/document-intake/failures").glob("*.json"))
    assert len(failures) == 1
    failure = json.loads(failures[0].read_text(encoding="utf-8"))
    assert failure["error"]["message"] == "Docling is not installed"


@pytest.mark.parametrize("suffix", ["docx", "xlsx", "pptx"])
def test_docling_converts_office_documents_without_path_leakage(
    tmp_path: Path, suffix: str
) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()
    source = tmp_path / f"sample.{suffix}"
    if suffix == "docx":
        from docx import Document

        document = Document()
        document.add_heading("DOCX 검증", 0)
        document.add_paragraph("로컬 수집")
        document.save(str(source))
    elif suffix == "xlsx":
        from openpyxl import Workbook  # type: ignore[import-untyped]

        workbook = Workbook()
        sheet = workbook.active
        sheet.append(["키", "값"])
        sheet.append(["source", "local"])
        workbook.save(source)
    else:
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "PPTX 검증"
        slide.placeholders[1].text = "로컬 수집"
        presentation.save(str(source))

    result = ingest_document_candidate(
        source, vault, source_locator=f"approved-drop/sample.{suffix}"
    )
    assert result.candidate is not None
    candidate = vault / result.candidate
    combined = "\n".join(
        (candidate / name).read_text(encoding="utf-8")
        for name in (
            "document.json",
            "document.raw.md",
            "document.md",
            "quality.json",
            "chunks.jsonl",
            "receipt.json",
        )
    )

    assert result.status == "candidate"
    assert str(tmp_path) not in combined
    assert str(Path.home()) not in combined


def test_docling_converts_scanned_pdf_with_local_rapidocr_cache(tmp_path: Path) -> None:
    cache_value = os.environ.get("WOON_DOCLING_MODEL_CACHE")
    if cache_value is None:
        pytest.skip("WOON_DOCLING_MODEL_CACHE is required for the PDF/OCR integration gate")

    from PIL import Image, ImageDraw, ImageFont

    cache = Path(cache_value).expanduser().resolve()
    vault = tmp_path / "vault"
    vault.mkdir()
    source = tmp_path / "scan.pdf"
    image = Image.new("RGB", (1600, 900), "white")
    draw = ImageDraw.Draw(image)
    font: Any
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 72)
    except OSError:
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 72)
        except OSError:
            font = ImageFont.load_default()
    draw.text((120, 180), "WOON DOCLING OCR", fill="black", font=font)
    draw.text((120, 320), "LOCAL ONLY TEST 2026", fill="black", font=font)
    image.save(source, "PDF", resolution=150.0)

    first = ingest_document_candidate(
        source,
        vault,
        source_locator="integration/scan.pdf",
        ocr="rapidocr",
        model_cache=cache,
    )
    replay = ingest_document_candidate(
        source,
        vault,
        source_locator="integration/scan.pdf",
        ocr="rapidocr",
        model_cache=cache,
    )
    assert first.candidate is not None
    candidate = vault / first.candidate
    markdown = (candidate / "document.md").read_text(encoding="utf-8")
    combined = "\n".join(
        (candidate / name).read_text(encoding="utf-8")
        for name in (
            "document.json",
            "document.raw.md",
            "document.md",
            "quality.json",
            "chunks.jsonl",
            "receipt.json",
        )
    )

    assert "WOON DOCLING OCR" in markdown
    assert "LOCAL ONLY TEST 2026" in markdown
    receipt = json.loads((candidate / "receipt.json").read_text(encoding="utf-8"))
    assert receipt["converter"]["packages"]["onnxruntime"]
    assert receipt["converter"]["packages"]["rapidocr"]
    assert replay.candidate_id == first.candidate_id
    assert replay.replayed is True
    assert str(tmp_path) not in combined
    assert str(cache) not in combined
    assert str(Path.home()) not in combined


def test_existing_source_catalog_classifies_docling_office_inputs_as_documents(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source"
    vault = tmp_path / "vault"
    source.mkdir()
    vault.mkdir()
    for name in ("sample.docx", "sample.pdf", "sample.pptx", "sample.xlsx"):
        (source / name).write_bytes(name.encode())

    plan = plan_source_catalog(source, vault, "document-drop")

    assert {record.locator: record.role for record in plan.records} == {
        "sample.docx": "document",
        "sample.pdf": "document",
        "sample.pptx": "document",
        "sample.xlsx": "document",
    }


def test_document_intake_cli_reports_candidate(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()
    output = StringIO()

    run(
        [
            "knowledge",
            "document-intake",
            "--source",
            str(FIXTURE),
            "--source-locator",
            "fixtures/table.html",
            "--vault",
            str(vault),
        ],
        output,
    )

    payload = json.loads(output.getvalue())
    assert payload["status"] == "candidate"
    assert payload["replayed"] is False
    assert payload["candidate"].startswith(
        ".local/woon-knowledge/document-intake/candidates/docling-"
    )
    assert payload["observation"].startswith(
        ".local/woon-knowledge/document-intake/observations/docling-"
    )


def test_discarded_document_candidate_is_terminal_without_visible_archive(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()
    candidate = ingest_document_candidate(FIXTURE, vault)
    decision = tmp_path / "decision.json"
    decision.write_text(
        json.dumps(
            {
                "version": 1,
                "candidate_id": candidate.candidate_id,
                "disposition": "discarded",
                "reason_code": "demo-only",
                "rationale": "변환 검증용 자료라 정본 지식으로 보존할 가치가 없습니다.",
                "canonical_paths": [],
                "source_targets": [],
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )

    result = resolve_document_candidate(vault, decision)
    replay = resolve_document_candidate(vault, decision)
    audit = audit_document_resolutions(vault)

    assert result.disposition == "discarded"
    assert result.replayed is False
    assert replay.replayed is True
    assert audit.complete is True
    assert audit.resolved == 1
    assert audit.pending == ()
    assert not (vault / "wiki").exists()


def test_legacy_review_candidate_can_only_be_discarded(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()
    candidate = ingest_document_candidate(FIXTURE, vault)
    directory = vault / str(candidate.candidate)
    receipt_path = directory / "receipt.json"
    receipt = json.loads(receipt_path.read_text(encoding="utf-8"))
    receipt["version"] = 1
    receipt["promotion_state"] = "review-required"
    receipt["outputs"] = {
        name: receipt["outputs"][name] for name in ("chunks.jsonl", "document.json", "document.md")
    }
    receipt_path.write_text(json.dumps(receipt), encoding="utf-8")
    for name in ("document.raw.md", "promotion.json", "quality.json"):
        (directory / name).unlink()
    decision = tmp_path / "legacy-discard.json"
    decision.write_text(
        json.dumps(
            {
                "version": 1,
                "candidate_id": candidate.candidate_id,
                "disposition": "discarded",
                "reason_code": "demo-only",
                "rationale": "이전 adapter 검증용 후보를 visible 문서 없이 종결합니다.",
                "canonical_paths": [],
                "source_targets": [],
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )

    result = resolve_document_candidate(vault, decision)

    assert result.disposition == "discarded"
    assert audit_document_resolutions(vault).complete is True


def test_integrated_document_requires_matching_wiki_owned_source(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()
    candidate = ingest_document_candidate(FIXTURE, vault)
    canonical = vault / "wiki/concepts/document-intake.md"
    canonical.parent.mkdir(parents=True)
    canonical.write_text("# 문서 수집\n", encoding="utf-8")
    source_target = vault / "wiki/private/_sources/knowledge/demo/table.html"
    source_target.parent.mkdir(parents=True)
    source_target.write_bytes(FIXTURE.read_bytes())
    decision = tmp_path / "integrated.json"
    decision.write_text(
        json.dumps(
            {
                "version": 1,
                "candidate_id": candidate.candidate_id,
                "disposition": "integrated",
                "reason_code": "existing-subject-updated",
                "rationale": "기존 문서 수집 주제에 검증 가능한 표 구조를 병합했습니다.",
                "canonical_paths": ["wiki/concepts/document-intake.md"],
                "source_targets": ["wiki/private/_sources/knowledge/demo/table.html"],
            },
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )

    result = resolve_document_candidate(vault, decision)

    assert result.disposition == "integrated"
    assert audit_document_resolutions(vault).complete is True


def test_document_resolution_audit_reports_unresolved_candidate(tmp_path: Path) -> None:
    vault = tmp_path / "vault"
    vault.mkdir()
    candidate = ingest_document_candidate(FIXTURE, vault)

    audit = audit_document_resolutions(vault)

    assert audit.complete is False
    assert audit.pending == (candidate.candidate_id,)
